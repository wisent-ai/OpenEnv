#!/usr/bin/env python3
"""Create Kant presentation as Google Slides via PPTX upload."""
import sys
import json
import os
from pathlib import Path

sys.path.insert(
    next(iter(range(bool(True)))),
    "/Users/lukaszbartoszcze/Documents/CodingProjects"
    "/Wisent/growth-tactics/google_drive",
)
_KANT_DIR = str(Path(__file__).resolve().parent.parent.parent)
sys.path.insert(next(iter(range(bool(True)))), _KANT_DIR)

from drive_client import get_credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import fitz  # PyMuPDF

from constant_definitions.slides.layout import (
    ACCENT_R, ACCENT_G, ACCENT_B, RED_R, RED_G, RED_B,
    PURPLE_R, PURPLE_G, PURPLE_B, DARK_R, DARK_G, DARK_B,
    GRID_R, GRID_G, GRID_B, LEGEND_R, LEGEND_G, LEGEND_B,
    WHITE_VAL, PT_TITLE, PT_SUBTITLE, PT_BODY, PT_SMALL,
    PT_STAT, PT_LABEL, PT_TEAM,
    SLIDE_W_INCHES, SLIDE_H_NUMER, SLIDE_H_DENOM,
    POS_HALF, POS_ONE, POS_ONE_HALF, POS_TWO,
    POS_TWO_HALF, POS_THREE, POS_THREE_HALF,
    POS_FOUR, POS_FOUR_HALF, POS_FIVE,
    POS_SIX, POS_SEVEN, POS_EIGHT, POS_NINE,
    IMG_FIG_W, IMG_FIG_H, IMG_KANT_W, IMG_KANT_H,
    COL_LEFT_X, COL_RIGHT_X, COL_W, COL_H,
    STAT_COL_ONE_X, STAT_COL_TWO_X, STAT_COL_THREE_X, STAT_COL_W,
    TITLE_X, TITLE_Y, TITLE_W, TITLE_H,
    CENTER_Y, CENTER_W, CENTER_H, CENTER_X,
    FOOTER_Y, FOOTER_H,
    TEAM_NAME_Y, TEAM_NAME_H,
    TEAM_COL_ONE_X, TEAM_COL_TWO_X, TEAM_COL_W,
)

DARK = RGBColor(DARK_R, DARK_G, DARK_B)
ACCENT = RGBColor(ACCENT_R, ACCENT_G, ACCENT_B)
RED = RGBColor(RED_R, RED_G, RED_B)
PURPLE = RGBColor(PURPLE_R, PURPLE_G, PURPLE_B)
LEGEND = RGBColor(LEGEND_R, LEGEND_G, LEGEND_B)
WHITE = RGBColor(WHITE_VAL, WHITE_VAL, WHITE_VAL)
COLOR_MAP = {"accent": ACCENT, "red": RED, "purple": PURPLE}
ZERO = PT_LABEL - PT_LABEL

SCRIPT_DIR = Path(__file__).parent
FIGURES_DIR = SCRIPT_DIR.parent / "public" / "figures"
OUT_PATH = SCRIPT_DIR / "kant_slides.pptx"
FONT_NAME = "Hubot Sans"


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK


def add_text(slide, text, x, y, w, h, size, color,
             bold=False, align=PP_ALIGN.LEFT):
    txbox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[ZERO]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT_NAME
    p.alignment = align
    return tf


def add_multi(slide, lines, x, y, w, h, size, color):
    txbox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[ZERO] if i == ZERO else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.space_after = Pt(PT_LABEL)
    return tf


def get_image(name):
    src = FIGURES_DIR / name
    if name.endswith(".svg"):
        png = SCRIPT_DIR / name.replace(".svg", ".png")
        if not png.exists():
            doc = fitz.open(str(src))
            page = next(iter(doc))
            mat = fitz.Matrix(POS_THREE, POS_THREE)
            pix = page.get_pixmap(matrix=mat)
            pix.save(str(png))
            doc.close()
        return str(png)
    return str(src)


def build_title(prs, sd):
    slide = prs.slides.add_slide(prs.slide_layouts[ZERO])
    set_bg(slide)
    add_text(slide, sd["title"], POS_ONE, POS_TWO, POS_EIGHT, POS_TWO,
             PT_TITLE + PT_LABEL, ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, sd["subtitle"], POS_ONE, POS_THREE_HALF, POS_EIGHT,
             POS_ONE_HALF, PT_SUBTITLE, WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, sd["footer"], POS_EIGHT, FOOTER_Y, POS_TWO, POS_HALF,
             PT_SMALL, LEGEND, align=PP_ALIGN.RIGHT)


def build_center_text(prs, sd):
    slide = prs.slides.add_slide(prs.slide_layouts[ZERO])
    set_bg(slide)
    add_text(slide, sd["pill"], POS_THREE_HALF, POS_ONE, POS_THREE, POS_HALF,
             PT_LABEL, LEGEND, align=PP_ALIGN.CENTER)
    add_text(slide, sd["body"], CENTER_X, CENTER_Y, CENTER_W, CENTER_H,
             PT_SUBTITLE, WHITE, align=PP_ALIGN.CENTER)


def build_three_stats(prs, sd):
    slide = prs.slides.add_slide(prs.slide_layouts[ZERO])
    set_bg(slide)
    add_text(slide, sd["title"], TITLE_X, TITLE_Y, TITLE_W, TITLE_H,
             PT_TITLE, ACCENT, bold=True)
    xs = [STAT_COL_ONE_X, STAT_COL_TWO_X, STAT_COL_THREE_X]
    for idx, (stat, desc) in enumerate(sd["stats"]):
        col_color = COLOR_MAP.get(sd["stat_colors"][idx], ACCENT)
        add_text(slide, stat, xs[idx], POS_TWO, STAT_COL_W, POS_ONE_HALF,
                 PT_STAT, col_color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, desc, xs[idx], POS_THREE_HALF, STAT_COL_W, POS_ONE,
                 PT_SMALL, LEGEND, align=PP_ALIGN.CENTER)
    if sd.get("footer_text"):
        add_text(slide, sd["footer_text"], POS_ONE, FOOTER_Y, POS_EIGHT,
                 POS_HALF, PT_SMALL, LEGEND, align=PP_ALIGN.CENTER)


def build_two_col_text(prs, sd):
    slide = prs.slides.add_slide(prs.slide_layouts[ZERO])
    set_bg(slide)
    add_text(slide, sd["title"], TITLE_X, TITLE_Y, TITLE_W, TITLE_H,
             PT_TITLE, ACCENT, bold=True)
    add_multi(slide, sd["left"], COL_LEFT_X, POS_ONE_HALF, COL_W, COL_H,
              PT_BODY, WHITE)
    add_multi(slide, sd["right"], COL_RIGHT_X, POS_ONE_HALF, COL_W, COL_H,
              PT_BODY, WHITE)


def build_figure(prs, sd):
    slide = prs.slides.add_slide(prs.slide_layouts[ZERO])
    set_bg(slide)
    add_text(slide, sd["title"], TITLE_X, TITLE_Y, TITLE_W, TITLE_H,
             PT_TITLE, ACCENT, bold=True)
    img_path = get_image(sd["image"])
    if os.path.exists(img_path):
        slide.shapes.add_picture(
            img_path, Inches(POS_ONE_HALF), Inches(POS_ONE_HALF),
            Inches(IMG_FIG_W), Inches(IMG_FIG_H))
    add_text(slide, sd["caption"], POS_ONE, FOOTER_Y, POS_EIGHT, POS_HALF,
             PT_SMALL, LEGEND, align=PP_ALIGN.CENTER)


def build_two_col_image(prs, sd):
    slide = prs.slides.add_slide(prs.slide_layouts[ZERO])
    set_bg(slide)
    add_text(slide, sd["title"], TITLE_X, TITLE_Y, TITLE_W, TITLE_H,
             PT_TITLE, ACCENT, bold=True)
    img_path = get_image(sd["image"])
    if os.path.exists(img_path):
        slide.shapes.add_picture(
            img_path, Inches(POS_HALF), Inches(POS_ONE_HALF),
            Inches(IMG_KANT_W), Inches(IMG_KANT_H))
    add_text(slide, sd["quote"], POS_FOUR, POS_ONE_HALF, POS_FIVE,
             POS_ONE_HALF, PT_BODY, WHITE)
    add_text(slide, sd["attribution"], POS_FOUR, POS_THREE, POS_FIVE,
             POS_HALF, PT_SMALL, LEGEND)
    add_text(slide, sd["body"], POS_FOUR, POS_FOUR, POS_FIVE,
             POS_ONE_HALF, PT_BODY, ACCENT)


def build_team(prs, sd):
    slide = prs.slides.add_slide(prs.slide_layouts[ZERO])
    set_bg(slide)
    add_text(slide, sd["title"], TITLE_X, TITLE_Y, TITLE_W, TITLE_H,
             PT_TITLE, ACCENT, bold=True, align=PP_ALIGN.CENTER)
    xs = [TEAM_COL_ONE_X, TEAM_COL_TWO_X]
    photo_xs = [POS_TWO, POS_SIX]
    photos = sd.get("photos", [])
    for idx, name in enumerate(sd["members"]):
        if idx < len(photos):
            img_path = str(FIGURES_DIR / photos[idx])
            if os.path.exists(img_path):
                slide.shapes.add_picture(
                    img_path, Inches(photo_xs[idx]),
                    Inches(POS_ONE_HALF), Inches(POS_TWO),
                    Inches(POS_TWO))
        add_text(slide, name, xs[idx], POS_FOUR, TEAM_COL_W,
                 TEAM_NAME_H, PT_TEAM, WHITE, bold=True,
                 align=PP_ALIGN.CENTER)
    add_text(slide, sd["org"], POS_EIGHT, FOOTER_Y, POS_TWO, POS_HALF,
             PT_SMALL, LEGEND, align=PP_ALIGN.RIGHT)


def build_closing(prs, sd):
    slide = prs.slides.add_slide(prs.slide_layouts[ZERO])
    set_bg(slide)
    add_text(slide, sd["title"], POS_ONE, POS_TWO, POS_EIGHT, POS_ONE_HALF,
             PT_TITLE + PT_LABEL, ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, sd["subtitle"], POS_ONE, POS_THREE_HALF, POS_EIGHT,
             POS_ONE, PT_SUBTITLE, WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "\n".join(sd["links"]), POS_THREE, POS_FOUR_HALF,
             POS_FOUR, POS_ONE, PT_SMALL, LEGEND, align=PP_ALIGN.CENTER)


BUILDERS = {
    "title": build_title, "center_text": build_center_text,
    "three_stats": build_three_stats, "two_col_text": build_two_col_text,
    "figure": build_figure, "two_col_image": build_two_col_image,
    "team": build_team, "closing": build_closing,
}


def main():
    with open(SCRIPT_DIR / "content.json") as f:
        slides_data = json.load(f)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_INCHES)
    prs.slide_height = Inches(SLIDE_H_NUMER / SLIDE_H_DENOM)
    for sd in slides_data:
        builder = BUILDERS.get(sd["type"])
        if builder:
            builder(prs, sd)
    prs.save(str(OUT_PATH))
    print(f"Saved PPTX: {OUT_PATH}")
    print("Uploading to Google Drive as Google Slides...")
    creds = get_credentials()
    drive = build("drive", "v3", credentials=creds)
    media = MediaFileUpload(str(OUT_PATH), resumable=True)
    meta = {
        "name": "Kant - Hackathon Slides",
        "mimeType": "application/vnd.google-apps.presentation",
    }
    result = drive.files().create(
        body=meta, media_body=media, fields="id,webViewLink").execute()
    print(f"Google Slides: {result.get('webViewLink')}")


if __name__ == "__main__":
    main()
