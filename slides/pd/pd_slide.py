#!/usr/bin/env python3
"""Generate a single Prisoner's Dilemma payoff-matrix slide and upload
to Google Slides. Uses the shared Wisent brand palette and layout
constants from constant_definitions.slides.layout.

This script produces a standalone PPTX with one slide containing a
styled payoff matrix table (header row plus action rows) along with a
brief explanation of the Nash equilibrium and alignment relevance.
"""
import sys
from pathlib import Path

sys.path.insert(
    next(iter(range(bool(True)))),
    "/Users/lukaszbartoszcze/Documents/CodingProjects"
    "/Wisent/growth-tactics/google_drive",
)
_KANT_DIR = str(Path(__file__).resolve().parent.parent.parent)
sys.path.insert(next(iter(range(bool(True)))), _KANT_DIR)

from drive_client import get_credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from constant_definitions.slides.layout import (
    ACCENT_R, ACCENT_G, ACCENT_B,
    DARK_R, DARK_G, DARK_B,
    LEGEND_R, LEGEND_G, LEGEND_B,
    WHITE_VAL, PT_TITLE, PT_BODY, PT_SMALL, PT_LABEL,
    SLIDE_W_INCHES, SLIDE_H_NUMER, SLIDE_H_DENOM,
    POS_HALF, POS_ONE, POS_ONE_HALF, POS_TWO,
    POS_THREE, POS_FOUR, POS_FIVE, POS_EIGHT,
    TITLE_X, TITLE_Y, TITLE_W, TITLE_H, FOOTER_Y,
    PD_CC, PD_CD, PD_DC, PD_DD,
    PD_NE_LABEL, PD_PO_LABEL, PD_EXPLANATION_BODY,
    PLAYER_ROW_LABEL, PLAYER_COL_LABEL,
)

DARK = RGBColor(DARK_R, DARK_G, DARK_B)
ACCENT = RGBColor(ACCENT_R, ACCENT_G, ACCENT_B)
LEGEND = RGBColor(LEGEND_R, LEGEND_G, LEGEND_B)
WHITE = RGBColor(WHITE_VAL, WHITE_VAL, WHITE_VAL)
GRID = RGBColor(
    DARK_R + PT_LABEL + PT_LABEL,
    DARK_G + PT_LABEL + PT_LABEL,
    DARK_B + PT_LABEL + PT_LABEL,
)
FONT_NAME = "Hubot Sans"

SCRIPT_DIR = Path(__file__).parent
OUT_PATH = SCRIPT_DIR / "pd_slide.pptx"

HEADER = ["", "Cooperate", "Defect"]
ROW_COOP = ["Cooperate", PD_CC, PD_CD]
ROW_DEF = ["Defect", PD_DC, PD_DD]
TABLE_ROWS = [HEADER, ROW_COOP, ROW_DEF]

EXPLANATION = PD_NE_LABEL + "\n" + PD_PO_LABEL + "\n\n" + PD_EXPLANATION_BODY

NUM_ROWS = len(TABLE_ROWS)
NUM_COLS = len(HEADER)
ZERO = NUM_ROWS - NUM_ROWS


def _cell_font(cell, text, size, color, bold=False, align=PP_ALIGN.CENTER):
    cell.text = ""
    p = cell.text_frame.paragraphs[ZERO]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT_NAME
    p.alignment = align
    cell.fill.solid()
    cell.fill.fore_color.rgb = GRID


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_INCHES)
    prs.slide_height = Inches(SLIDE_H_NUMER / SLIDE_H_DENOM)
    slide = prs.slides.add_slide(prs.slide_layouts[ZERO])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK

    # Title
    txbox = slide.shapes.add_textbox(
        Inches(TITLE_X), Inches(TITLE_Y),
        Inches(TITLE_W), Inches(TITLE_H))
    p = txbox.text_frame.paragraphs[ZERO]
    p.text = "Prisoner\u2019s Dilemma"
    p.font.size = Pt(PT_TITLE)
    p.font.color.rgb = ACCENT
    p.font.bold = True
    p.font.name = FONT_NAME

    # Column player label
    txbox = slide.shapes.add_textbox(
        Inches(POS_TWO + POS_ONE), Inches(POS_ONE),
        Inches(POS_FOUR), Inches(POS_HALF))
    p = txbox.text_frame.paragraphs[ZERO]
    p.text = PLAYER_COL_LABEL
    p.font.size = Pt(PT_BODY)
    p.font.color.rgb = LEGEND
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    # Row player label
    txbox = slide.shapes.add_textbox(
        Inches(POS_HALF), Inches(POS_TWO + POS_HALF),
        Inches(POS_ONE_HALF), Inches(POS_HALF))
    p = txbox.text_frame.paragraphs[ZERO]
    p.text = PLAYER_ROW_LABEL
    p.font.size = Pt(PT_BODY)
    p.font.color.rgb = LEGEND
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    # Payoff matrix table
    tbl = slide.shapes.add_table(
        NUM_ROWS, NUM_COLS,
        Inches(POS_TWO), Inches(POS_ONE_HALF),
        Inches(POS_FIVE), Inches(POS_TWO)).table

    for r_idx, row_data in enumerate(TABLE_ROWS):
        for c_idx, val in enumerate(row_data):
            cell = tbl.cell(r_idx, c_idx)
            is_header = (r_idx == ZERO)
            is_label = (c_idx == ZERO)
            color = ACCENT if (is_header or is_label) else WHITE
            _cell_font(cell, val, PT_BODY, color,
                       bold=(is_header or is_label))

    # Explanation text block
    txbox = slide.shapes.add_textbox(
        Inches(POS_HALF), Inches(POS_THREE + POS_ONE),
        Inches(POS_EIGHT + POS_ONE), Inches(POS_ONE_HALF))
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(EXPLANATION.split("\n")):
        para = tf.paragraphs[ZERO] if i == ZERO else tf.add_paragraph()
        para.text = line
        para.font.size = Pt(PT_SMALL)
        para.font.color.rgb = LEGEND
        para.font.name = FONT_NAME

    # Footer
    txbox = slide.shapes.add_textbox(
        Inches(POS_EIGHT), Inches(FOOTER_Y),
        Inches(POS_TWO), Inches(POS_HALF))
    p = txbox.text_frame.paragraphs[ZERO]
    p.text = "Wisent"
    p.font.size = Pt(PT_SMALL)
    p.font.color.rgb = LEGEND
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.RIGHT

    prs.save(str(OUT_PATH))
    print(f"Saved PPTX: {OUT_PATH}")

    print("Uploading to Google Drive as Google Slides...")
    creds = get_credentials()
    drive = build("drive", "v3", credentials=creds)
    media = MediaFileUpload(str(OUT_PATH), resumable=True)
    meta = {
        "name": "Kant - Prisoner\u2019s Dilemma",
        "mimeType": "application/vnd.google-apps.presentation",
    }
    result = drive.files().create(
        body=meta, media_body=media, fields="id,webViewLink").execute()
    print(f"Google Slides: {result.get('webViewLink')}")


if __name__ == "__main__":
    main()
